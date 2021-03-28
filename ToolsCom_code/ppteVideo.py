import sys, getopt
import cv2
from pptx import Presentation
from pptx.util import Inches
from os.path import exists
class ptvcom():
	def startd(self,videoFindr,powerNamev,powerFindv,videoStartk,videoEndk,videoA,videoB):
	    powerK = Presentation()
	    for i in range(videoStartk,videoEndk):
	        zw,videoR = videoFindr.read()
	        cv2.imwrite(powerFindv + '/x.png',videoR)
	        slide = powerK.slides.add_slide(powerK.slide_layouts[6])
	        img_path = powerFindv + '/x.png'
	        pic = slide.shapes.add_picture(img_path,Inches(-1), Inches(0), Inches(12), Inches(6.75))
	        powerK.save(powerFindv + '/' + powerNamev + '.pptx')
	        videoProgressbar['value'] = i * (100 / (videoEndk - videoStartk))
	def powerVideo(argv):
		if len(sys.argv) != 8:
			print('参数至少输入7个\n需要参数\n视频名称、ppt路径、ppt名称、视频起始帧数、视频结束帧数、视频长宽')
		videoX = sys.argv
		if exists(videoX[1]) == False or exists(videoX[2]) == False:
			print('路径无效或路径不可用')
		elif videoX[4]  >= videoX[5]:
			print('视频起始帧数不能大于或等于视频结束帧数')
		else:
			startd(videoX[1],videoX[3],videoX[2],videoX[4],videoX[5],videoX[6],videoX[7])
	def in_it():
		ptvcom.powerVideo(sys.argv[1:])


ptvcom.in_it()
