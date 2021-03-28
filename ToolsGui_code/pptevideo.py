###引入库部分###pror
# -*- coding: UTF-8 -*-
import tkinter as tk
import tkinter.ttk
import tkinter.filedialog
import tkinter.messagebox
import time
import cv2
from pptx import Presentation
from pptx.util import Inches
from os.path import exists
window = tk.Tk()
window.title('PpteVideo')
window.geometry('1000x700')
###引入库部分###

###逻辑代码部分###
#视频转ppt
def Vixpower(frameNum,frameSecond,videoFindr,powerNamev,powerFindv,videoStartk,videoEndk):
    print(frameNum,frameSecond,videoFindr,powerNamev,powerFindv)
    powerK = Presentation()
    for i in range(videoStartk,videoEndk):
        print(powerFindv + '/x.png')
        zw,videoR = videoFindr.read()
        cv2.imwrite(powerFindv + '/x.png',videoR)
        slide = powerK.slides.add_slide(powerK.slide_layouts[6])
        img_path = powerFindv + '/x.png'
        pic = slide.shapes.add_picture(img_path,Inches(-1), Inches(0), Inches(12), Inches(6.75))
        powerK.save(powerFindv + '/' + powerNamev + '.pptx')
        videoProgressbar['value'] = i * (100 / (videoEndk - videoStartk))
        window.update()
    videoProgressbar['value'] = 100
    window.update()
#检测变量错误，并将参数传递给Vixpower
def tryStartVideoExPower():
	if VideoNameVar.get() == '':
		tk.messagebox.showinfo('提示','您视频丢了？')
	elif exists(VideoNameVar.get()) == False:
		tk.messagebox.showinfo('提示','电脑里没有这个文件')
	elif powerName.get() == '':
		tkinter.messagebox.askquestion('提示','您要用默认输出名称？')
	elif PowerFindVar.get() == '':
		tkinter.messagebox.showinfo('提示','ppt保存在哪！？')
	elif exists(VideoNameVar.get()) == False:
		tkinter.messagebox.showinfo('提示','哪有这个路径啊？？？')
	elif videoStartVar.get() >= videoEndVar.get():
		tkinter.messagebox.showinfo('提示','结束位置不能小于或等于起始位置！！！')
	else:
		Vixpower(cv2.VideoCapture(VideoNameVar.get()).get(7), #opencv模块读取视频名称的总帧数
				 cv2.VideoCapture(VideoNameVar.get()).get(5), #opencv模块读取视频名称的帧率
				 cv2.VideoCapture(VideoNameVar.get()),
				 powerName.get(),
				 PowerFindVar.get(),
				 videoStartVar.get(),
				 videoEndVar.get())
#修改滑动条的状态
def getVideoFrame(fall):
    videoRead = cv2.VideoCapture(fall)
    videoFrameSecond = videoRead.get(7)
    videoReadAll = videoRead.get(5)
    videoStartScale['to'] = videoFrameSecond
    videoStartScale['tickinterval'] = videoFrameSecond/15
    videoStartScale['state'] = 'normal'
    videoEndScale['to'] = videoFrameSecond
    videoEndScale['tickinterval'] = videoFrameSecond/15
    videoEndScale['state'] = 'normal'
#读取视频
def VideoNameFind():
    tineVar = tk.filedialog.askopenfilename(filetypes=[("视频",".mp4")])
    VideoNameVar.set(tineVar)
    getVideoFrame(fall = VideoNameVar.get())
#读取ppt路径
def powerNameFind():
    tineVar = tk.filedialog.askdirectory()
    PowerFindVar.set(tineVar)
###逻辑代码部分###

###GUI部分###
tk.Label(window, text='PpteVideo', font=('Arial', 35), width=30, height=2).pack()
#定义标题
videoFinder = tk.Frame(window)
videoFinder.pack()
VideoNameVar = tk.Variable()
tk.Label(videoFinder, text='视频路径:', font=(None,20)).pack(side='left')
tk.Entry(videoFinder, show=None, font=(None, 19),width=40, textvariable=VideoNameVar).pack(side='left',fill='x')
tk.Button(videoFinder, text='路径选择',font=(None, 14), command=VideoNameFind).pack(expand='YES',side='left')
#路径选择的frame
powerFinder = tk.Frame(window)
powerFinder.pack()
powerName = tk.Variable()
tk.Label(powerFinder, text='输出名称:', font=(None,20)).pack(side='left')
tk.Entry(powerFinder, show=None, font=(None, 19),width=40, textvariable=powerName).pack(side='left',fill='x')
tk.Label(powerFinder, text='.pptx ', font=(None,20)).pack(side='left')
#ppt输出的frame
videoEndFinder = tk.Frame(window)
videoEndFinder.pack()
PowerFindVar = tk.Variable()
tk.Label(videoEndFinder, text='输出路径:', font=(None,20)).pack(side='left')
tk.Entry(videoEndFinder, show=None, font=(None, 19),width=40, textvariable=PowerFindVar).pack(side='left',fill='x')
tk.Button(videoEndFinder, text='路径选择',font=(None, 14), command=powerNameFind).pack(expand='YES',side='left')
#输出路径的frame
videoStartFPS = tk.Frame(window)
videoStartFPS.pack()
videoStartVar = tk.Variable()
tk.Label(videoStartFPS, text='\n视频起始位置(帧数)', font=(None,20)).pack(side='top')
videoStartScale = tk.Scale(videoStartFPS,from_=0, to=210, orient=tk.HORIZONTAL, length=800, variable=videoStartVar, showvalue=1,tickinterval=30,bd=1,sliderrelief='flat',state="disabled")
videoStartScale.pack(side='top')
#视频起始帧数frame
videoEndFPS = tk.Frame(window)
videoEndFPS.pack()
videoEndVar = tk.Variable()
tk.Label(videoEndFPS, text='视频结束位置(帧数)', font=(None,20)).pack(side='top')
videoEndScale = tk.Scale(videoEndFPS,from_=0, to=210, orient=tk.HORIZONTAL, length=800, variable = videoEndVar, showvalue=1,tickinterval=30,bd=1,sliderrelief='flat',state="disabled")
videoEndScale.pack(side='top')
#视频结束帧数frame

tk.Label(window, text='  ', font=('Arial', 20), width=30, height=2).pack()

tk.Button(window, text='导出', font=(None, 25), width=20, height=1, command= tryStartVideoExPower).pack()

tk.Label(window, text='\n导出进度', font=(None, 15), width=30, height=2).pack()
videoProgressbar = tkinter.ttk.Progressbar(window, length=800)
videoProgressbar.pack()
window.mainloop()
#导出按钮
###GUI部分###
